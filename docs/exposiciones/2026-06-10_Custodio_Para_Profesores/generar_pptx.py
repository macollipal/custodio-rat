"""
Genera PPTX educativo: Custodio RAT Manager para profesores de Ley 21.719.
10 slides, tono sobrio, resumido.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_FILE = os.path.join(OUT_DIR, "Custodio_Para_Profesores_Ley21719.pptx")

# --- Colores sobrios ---
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)   # Azul oscuro profesional
MID_BLUE  = RGBColor(0x2C, 0x5F, 0x8E)   # Azul medio
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)  # Fondo slide
TEXT_DARK  = RGBColor(0x1F, 0x29, 0x37)  # Texto principal
TEXT_MID   = RGBColor(0x4B, 0x55, 0x63)  # Texto secundario
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x25, 0x63, 0xEB) # Azul Vercel-like

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=TEXT_DARK,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def header_bar(slide, title):
    """Barra de título en la parte superior."""
    add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
    add_text_box(slide, title, 0.5, 0.18, 12, 0.6,
                 font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def bullet_block(slide, items, left, top, width, font_size=17, color=TEXT_DARK):
    """Lista de bullets."""
    y = top
    for item in items:
        add_text_box(slide, f"  \u2022  {item}", left, y, width, 0.45,
                     font_size=font_size, color=color)
        y += 0.52

def slide_num(slide, num):
    add_text_box(slide, str(num), 12.5, 7.0, 0.5, 0.3,
                 font_size=11, color=TEXT_MID, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 1 — PORTADA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), SLIDE_H / Inches(1), DARK_BLUE)
add_rect(slide, 0, 5.5, SLIDE_W / Inches(1), 2.0, MID_BLUE)

add_text_box(slide, "Custodio RAT Manager", 1.0, 1.8, 11, 1.2,
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "Gestión del Registro de Actividades de Tratamiento", 1.0, 3.1, 11, 0.7,
             font_size=22, color=RGBColor(0xBF, 0xD7, 0xF0), align=PP_ALIGN.CENTER)
add_text_box(slide, "Ley 21.719 — Chile", 1.0, 3.9, 11, 0.5,
             font_size=18, color=RGBColor(0x93, 0xC3, 0xF9), align=PP_ALIGN.CENTER)
add_text_box(slide, "Exposición para profesores de Derecho — 10 de junio de 2026",
             1.0, 5.7, 11, 0.4, font_size=14, color=RGBColor(0x93, 0xC3, 0xF9),
             align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — QUÉ ES CUSTODIO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
header_bar(slide, "\u00bfQu\u00e9 es Custodio?")
add_text_box(slide, "Custodio es un software de gesti\u00f3n del Registro de Actividades de Tratamiento (RAT) dise\u00f1ado para empresas chilenas que procesan datos personales.", 0.5, 1.1, 12.3, 0.9, font_size=17, color=TEXT_DARK)
add_rect(slide, 0.5, 2.1, 12.3, 0.04, ACCENT)
bullet_block(slide, [
    "Permite crear, editar y mantener actualizado el RAT de cada empresa",
    "Gestiona tickets de derechos ARCO (Acceso, Rectificaci\u00f3n, Cancelaci\u00f3n, Oposici\u00f3n)",
    "Registra brechas de seguridad y notificaci\u00f3nes a la APDP",
    "Incluye m\u00f3dulos de transparencia, contratos de encargados y consentimientos",
    "Genera reportes y exporta PDF para auditor\u00edas",
    "Automatiza el cumplimiento de la Ley 21.719 (Art. 15 y mods.)",
], 0.5, 2.25, 12.3, font_size=16, color=TEXT_DARK)
slide_num(slide, 2)

# ============================================================
# SLIDE 3 — MARCO LEGAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
header_bar(slide, "Marco Legal — Ley 21.719")
bullet_block(slide, [
    "Ley 21.719 (feb. 2023):取代 la antigua Ley 19.628 sobre protecci\u00f3n de datos personales en Chile.",
    "Art. 15: Toda empresa que trate datos personales debe mantener un RAT actualizado.",
    "Art. 16: Datos sensibles requieren consentimiento expreso e informaci\u00f3n adicional.",
    "Art. 14 bis: Brechas de seguridad deben notificarse a la APDP en 72 horas.",
    "Art. 17: Los titulares tienen derechos ARCO que las empresas deben atender.",
    "Custodio centraliza todas estas obligaciones en una \u00fanica plataforma.",
], 0.5, 1.1, 12.3, font_size=16, color=TEXT_DARK)
add_rect(slide, 0.5, 5.3, 12.3, 0.8, RGBColor(0xE8, 0xF0, 0xFE))
add_text_box(slide, "Nota: La Ley 21.719 incorpora est\u00e1ndares del RGPD europeo y se aplica a todo tratamiento de datos personales realizado en Chile.", 0.7, 5.35, 11.9, 0.6, font_size=13, color=RGBColor(0x1E, 0x40, 0xAF))
slide_num(slide, 3)

# ============================================================
# SLIDE 4 — DASHBOARD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
header_bar(slide, "Dashboard — Panel Principal")
add_text_box(slide, "El dashboard es la pantalla de inicio tras el login. Muestra un resumen ejecutivo del estado de cumplimiento de la empresa.", 0.5, 1.1, 12.3, 0.7, font_size=16, color=TEXT_DARK)
add_rect(slide, 0.5, 1.85, 12.3, 0.04, ACCENT)
bullet_block(slide, [
    "KPIs principales: Total de RATs, RATs aprobados, en borrador, con datos sensibles.",
    "Gr\u00e1ficos de estado: RATs por estado (borrador, completo, en revisi\u00f3n, aprobado).",
    "Alertas: Brechas pendientes, tickets ARCO cerca del vencimiento SLA.",
    "Checklist de primeros pasos: Completar empresa, definir DPO, crear primer RAT, etc.",
    "Acceso r\u00e1pido a todas las secciones del men\u00fa lateral.",
], 0.5, 2.0, 12.3, font_size=16, color=TEXT_DARK)
slide_num(slide, 4)

# ============================================================
# SLIDE 5 — PROCESOS RAT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
header_bar(slide, "Procesos RAT — Registro de Actividades de Tratamiento")
add_text_box(slide, "El RAT es el coraz\u00f3n de la ley. Cada RAT representa un tratamiento espec\u00edfico de datos personales que realiza la empresa.", 0.5, 1.1, 12.3, 0.7, font_size=16, color=TEXT_DARK)
add_rect(slide, 0.5, 1.85, 12.3, 0.04, ACCENT)
bullet_block(slide, [
    "Nombre del proceso: Identificador \u00fanico (ej: \u201cGesti\u00f3n de n\u00f3mina\u201d).",
    "Categor\u00eda de titulares: Empleados, clientes, proveedores, etc.",
    "Categor\u00eda de datos: Datos de contacto, financieros, biom\u00e9tricos, etc.",
    "Finalidad: Para qu\u00e9 se usan los datos (ej: gesti\u00f3n de personal).",
    "Base legal: Consentimiento, contrato, obligaci\u00f3n legal o inter\u00e9s leg\u00edtimo.",
    "Plazo de retenci\u00f3n: Tiempo que se conservan los datos.",
], 0.5, 2.0, 12.3, font_size=16, color=TEXT_DARK)
add_rect(slide, 0.5, 5.3, 12.3, 0.8, RGBColor(0xE8, 0xF0, 0xFE))
add_text_box(slide, "Art. 15 Ley 21.719: El RAT debe contener al menos estos 9 campos. Custodio los valida autom\u00e1ticamente.", 0.7, 5.35, 11.9, 0.6, font_size=13, color=RGBColor(0x1E, 0x40, 0xAF))
slide_num(slide, 5)

# ============================================================
# SLIDE 6 — WIZARD RAT (4 PASOS)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
header_bar(slide, "Creaci\u00f3n de RAT — Wizard de 4 Pasos")
# 4 cajas de color para los pasos
steps = [
    ("1. Identificaci\u00f3n", "Nombre, categor\u00eda de titulares, fuente de datos, destinatarios."),
    ("2. Datos Tratados", "Categor\u00eda de datos, datos sensibles, EIPD, decisiones automatizadas."),
    ("3. Finalidad y Ley", "Finalidad, base legal, test de inter\u00e9s leg\u00edtimo (si aplica)."),
    ("4. Almacenamiento", "Plazo de retenci\u00f3n, medidas de seguridad, transferencias internacionales."),
]
x = 0.5
for title, desc in steps:
    box = add_rect(slide, x, 1.2, 2.9, 2.2, LIGHT_GRAY)
    box.line.color.rgb = ACCENT
    add_rect(slide, x, 1.2, 2.9, 0.5, ACCENT)
    add_text_box(slide, title, x + 0.1, 1.22, 2.7, 0.45,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.1, 1.78, 2.7, 1.5, font_size=13, color=TEXT_DARK)
    x += 3.1

add_text_box(slide, "El wizard gu\u00eda paso a paso al usuario, validando cada campo antes de avanzar. Al final, el RAT queda en estado \u201cborrador\u201d para revisi\u00f3n antes de aprobar.", 0.5, 3.6, 12.3, 0.7, font_size=15, color=TEXT_MID)
add_text_box(slide, "Estados del RAT: Borrador \u2192 Completo \u2192 En Revisi\u00f3n \u2192 Aprobado", 0.5, 4.4, 12.3, 0.5, font_size=15, bold=True, color=MID_BLUE)
slide_num(slide, 6)

# ============================================================
# SLIDE 7 — TICKETS ARCO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
header_bar(slide, "Tickets ARCO — Derechos de los Titulares")
add_text_box(slide, "Cuando un titular ejerci\u00f3 un derecho ARCO, la empresa debe registrar y gestionar el ticket en Custodio.", 0.5, 1.1, 12.3, 0.7, font_size=16, color=TEXT_DARK)
add_rect(slide, 0.5, 1.85, 12.3, 0.04, ACCENT)
bullet_block(slide, [
    "Acceso: El titular solicita conocer qu\u00e9 datos suyos tiene la empresa.",
    "Rectificaci\u00f3n: Solicita corregir datos inexactos o incompletos.",
    "Cancelaci\u00f3n: Solicita eliminar sus datos (derecho al olvido).",
    "Oposici\u00f3n: Se opone a un tratamiento espec\u00edfico.",
    "Custodio registra cada ticket con estado (nuevo / en proceso / resuelto) y SLA.",
    "El sistema muestra un historial de cambios y la respuesta de la empresa.",
], 0.5, 2.0, 12.3, font_size=16, color=TEXT_DARK)
add_rect(slide, 0.5, 5.3, 12.3, 0.8, RGBColor(0xE8, 0xF0, 0xFE))
add_text_box(slide, "Art. 17 Ley 21.719: La empresa debe responder al titular en un plazo razonable. Custodio calcula el SLA autom\u00e1ticamente.", 0.7, 5.35, 11.9, 0.6, font_size=13, color=RGBColor(0x1E, 0x40, 0xAF))
slide_num(slide, 7)

# ============================================================
# SLIDE 8 — BRECHAS DE SEGURIDAD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
header_bar(slide, "Brechas de Seguridad — Art. 14 bis")
add_text_box(slide, "Una brecha de seguridad ocurre cuando datos personales se ven comprometidos por acceso no autorizado, p\u00e9rdida o destrucci\u00f3n.", 0.5, 1.1, 12.3, 0.7, font_size=16, color=TEXT_DARK)
add_rect(slide, 0.5, 1.85, 12.3, 0.04, ACCENT)
bullet_block(slide, [
    "Custodio permite registrar la brecha con fecha de detecci\u00f3n, descripci\u00f3n y nivel de riesgo.",
    "Niveles de riesgo: Bajo, Medio, Alto, Cr\u00edtico.",
    "Si hay datos sensibles o menores afectados, se debe notificar a los titulares.",
    "Notificaci\u00f3n obligatoria a la APDP dentro de 72 horas (Art. 14 bis).",
    "Custodio lleva control de si la brecha fue notificada y a qui\u00e9n.",
    "Se pueden asociar RATs afectados a la brecha para trazabilidad.",
], 0.5, 2.0, 12.3, font_size=16, color=TEXT_DARK)
add_rect(slide, 0.5, 5.3, 12.3, 0.8, RGBColor(0xFE, 0xF3, 0xC7))
add_text_box(slide, "Importante: Las 72 horas se cuentan desde la detecci\u00f3n, no desde el descubrimiento. La notificaci\u00f3n debe incluir nature, consecuencias y medidas adoptadas.", 0.7, 5.35, 11.9, 0.6, font_size=13, color=RGBColor(0x92, 0x40, 0x0D))
slide_num(slide, 8)

# ============================================================
# SLIDE 9 — MÓDULOS DE CUMPLIMIENTO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), 0.9, DARK_BLUE)
header_bar(slide, "M\u00f3dulos de Cumplimiento")
modules = [
    ("Transparencia\n(Art. 14 ter)", "Pol\u00edtica de transparencia p\u00fablica que la empresa debe publicar."),
    ("Enc. Contrato\n(Art. 14 quater)", "Contratos con terceros que trattan datos por cuenta de la empresa."),
    ("Consentimientos\n(Art. 12)", "Registro de consentimientos v\u00e1lidos de los titulares."),
    ("EIPD\n(Art. 15 bis)", "Evaluaci\u00f3n de Impacto en Protecci\u00f3n de Datos (obligatoria para datos sensibles)."),
]
x = 0.5
for title, desc in modules:
    box = add_rect(slide, x, 1.15, 2.9, 2.5, LIGHT_GRAY)
    box.line.color.rgb = ACCENT
    add_rect(slide, x, 1.15, 2.9, 0.55, ACCENT)
    add_text_box(slide, title, x + 0.1, 1.17, 2.7, 0.5,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.1, 1.77, 2.7, 1.7, font_size=13, color=TEXT_DARK)
    x += 3.1

add_text_box(slide, "Estos m\u00f3dulos complementan el RAT y aseguran que la empresa cumpla todas las obligaciones de la Ley 21.719 de forma centralizada.", 0.5, 3.85, 12.3, 0.65, font_size=15, color=TEXT_MID)
slide_num(slide, 9)

# ============================================================
# SLIDE 10 — CIERRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SLIDE_W / Inches(1), SLIDE_H / Inches(1), DARK_BLUE)
add_rect(slide, 0, 5.5, SLIDE_W / Inches(1), 2.0, MID_BLUE)
add_text_box(slide, "Custodio: Cumplimiento simplificado", 1.0, 1.5, 11, 0.9,
             font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 4.0, 2.55, 5.3, 0.04, ACCENT)
bullet_block(slide, [
    "RAT completo y actualizado para auditor\u00edas.",
    "Gesti\u00f3n de tickets ARCO con SLA autom\u00e1tico.",
    "Registro de brechas con notificaci\u00f3n a la APDP.",
    "M\u00f3dulos de transparencia, contratos, consentimientos y EIPD.",
    "Reportes y exportaci\u00f3n PDF para evidencia de cumplimiento.",
    "Acceso desde cualquier navegador — sin instalaci\u00f3n.",
], 1.0, 2.75, 11, font_size=17, color=RGBColor(0xBF, 0xD7, 0xF0))
add_text_box(slide, "Para m\u00e1s informaci\u00f3n: www.custodio.cl  |  Versi\u00f3n beta en evaluaci\u00f3n",
             1.0, 5.7, 11, 0.5, font_size=15, color=RGBColor(0x93, 0xC3, 0xF9),
             align=PP_ALIGN.CENTER)
slide_num(slide, 10)

prs.save(OUT_FILE)
print(f"PPT guardado en: {OUT_FILE}")